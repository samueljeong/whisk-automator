#!/usr/bin/env python3
"""토스페이먼츠 카드사 심사용 결제경로 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SCREENSHOTS = {
    "product": "/Users/samuel/Downloads/debug-reports/mcp-screenshot-2026-03-11T14-07-29.png",
    "refund_footer": "/Users/samuel/Downloads/debug-reports/mcp-screenshot-2026-03-11T14-09-45.png",
    "billing": "/Users/samuel/Downloads/debug-reports/mcp-screenshot-2026-03-11T14-10-15.png",
}

OUTPUT = "/Users/samuel/Projects/코딩/whisk-automator/payment/결제경로_FlowGrokPro.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, title, subtitle=None):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(8)


def add_image_centered(slide, img_path, top=Inches(1.5), max_w=Inches(11.5), max_h=Inches(5.5)):
    from PIL import Image
    img = Image.open(img_path)
    w, h = img.size
    ratio = w / h

    if max_w / max_h > ratio:
        height = max_h
        width = int(height * ratio)
    else:
        width = max_w
        height = int(width / ratio)

    left = (prs.slide_width - width) // 2
    slide.shapes.add_picture(img_path, left, top, width, height)


# ── Slide 1: 가맹점 정보 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
add_title_text(slide, "1. 가맹점 정보")

info_lines = [
    ("상호", "홈즈컴퍼니"),
    ("대표", "서연희"),
    ("사업자등록번호", "467-04-01758"),
    ("통신판매업신고", "제 2020-의정부신곡-0232 호"),
    ("주소", "경기도 의정부시 장금로78번길 41, 102동 907호"),
    ("연락처", "010-8673-4129"),
    ("서비스 URL", "https://payment-tau-three.vercel.app"),
    ("서비스명", "Flow+Grok Pro (크롬 확장프로그램)"),
    ("결제 방식", "정기결제 (빌링) — 월간 ₩9,900 / 연간 ₩100,000"),
]

y = Inches(1.8)
for label, value in info_lines:
    txBox = slide.shapes.add_textbox(Inches(1.5), y, Inches(10), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{label}:  "
    run1.font.size = Pt(18)
    run1.font.bold = True
    run1.font.color.rgb = PURPLE
    run2 = p.add_run()
    run2.text = value
    run2.font.size = Pt(18)
    run2.font.color.rgb = WHITE
    y += Inches(0.55)


# ── Slide 2: 하단 사업자정보 (footer) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_text(slide, "2. 하단 사업자정보", "결제 페이지 하단 — 사업자 정보 표시")
add_image_centered(slide, SCREENSHOTS["refund_footer"])


# ── Slide 3: 환불규정 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_text(slide, "3. 교환 및 환불 규정", "URL: payment-tau-three.vercel.app/?demo=true  →  하단 스크롤")
add_image_centered(slide, SCREENSHOTS["refund_footer"])


# ── Slide 4: 상품 선택 / 구매 과정 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_text(slide, "4. 상품 선택 및 구매 과정", "월간(₩9,900) / 연간(₩100,000) 플랜 선택 → '구독하기' 클릭")
add_image_centered(slide, SCREENSHOTS["product"])


# ── Slide 5: 빌링 카드 입력창 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_text(slide, "5. 토스페이먼츠 빌링 카드 등록", "'구독하기' 클릭 후 → 토스페이먼츠 카드 입력 창 표시")
add_image_centered(slide, SCREENSHOTS["billing"])


prs.save(OUTPUT)
print(f"PPT 생성 완료: {OUTPUT}")
